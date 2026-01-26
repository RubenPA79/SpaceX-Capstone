from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    # Load template if exists, else create new
    template_path = "ds-capstone-template-coursera (1).pptx"
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"Loaded template: {template_path}")
    else:
        prs = Presentation()
        print("Template not found, created new presentation.")

    # Helper to add slide
    def add_slide(title, content_text=None, image_path=None):
        # Use a blank layout or title+content
        # Layouts: 0=Title, 1=Title+Content, ...
        # We'll try to use Title+Content (index 1 usually)
        if len(prs.slide_layouts) > 1:
            slide_layout = prs.slide_layouts[1] 
        else:
            slide_layout = prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Set Content (Text)
        if content_text:
            # Find the placeholder for text
            # Usually the second placeholder
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = content_text
            else:
                # Add text box if no placeholder
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = content_text
                
        # Add Image
        if image_path:
            if os.path.exists(image_path):
                # Add image on the right or bottom
                left = Inches(2)
                top = Inches(3.5)
                height = Inches(3.5)
                slide.shapes.add_picture(image_path, left, top, height=height)
            else:
                print(f"Warning: Image not found {image_path}")

    # --- SLIDES ---

    # 1. Executive Summary
    add_slide(
        "Resumen Ejecutivo",
        "El objetivo es predecir el éxito del aterrizaje de la primera etapa del Falcon 9.\n"
        "Esto reduce costos de lanzamiento significativamente.\n"
        "Metodología: Recopilación de datos (API, Web Scraping), EDA, SQL, Mapas, Dash, Machine Learning.\n"
        "Resultado: El modelo de clasificación alcanzó una alta precisión (~83%)."
    )

    # 2. Introduction
    add_slide(
        "Introducción",
        "Contexto: SpaceX revoluciona la industria aeroespacial con cohetes reutilizables.\n"
        "Problema: Predecir si la primera etapa aterrizará exitosamente.\n"
        "Datos: Histórico de lanzamientos de Falcon 9."
    )

    # 3. Methodology
    add_slide(
        "Metodología",
        "1. Recopilación de Datos: API SpaceX, Wikipedia.\n"
        "2. Data Wrangling: Limpieza, manejo de nulos.\n"
        "3. EDA: Visualizaciones y SQL.\n"
        "4. Visualización Interactiva: Folium y Dash.\n"
        "5. Análisis Predictivo: Modelos de Clasificación (LR, SVM, Tree, KNN)."
    )

    # 4. EDA - Visuals
    add_slide(
        "Resultados EDA - Tendencia Temporal",
        "La tasa de éxito ha mejorado consistentemente a lo largo de los años, demostrando la madurez tecnológica de SpaceX.",
        "assets/eda/7_yearly_trend.png"
    )

    add_slide(
        "Resultados EDA - Carga vs Sitio",
        "Distribución de cargas útiles por sitio de lanzamiento. KSC LC-39A maneja cargas variadas con alto éxito.",
        "assets/eda/3_payload_vs_site.png"
    )

    add_slide(
        "Resultados EDA - Éxito por Órbita",
        "Las órbitas ES-L1, GEO, HEO y SSO muestran las tasas de éxito más altas.",
        "assets/eda/4_success_by_orbit.png"
    )

    # 5. SQL Analysis
    add_slide(
        "Resultados - Análisis SQL",
        "Consultas clave realizadas:\n"
        "- Identificación de nombres únicos de sitios de lanzamiento.\n"
        "- Cálculo de la carga útil total transportada por la NASA.\n"
        "- Análisis de resultados de aterrizaje por fecha.\n"
        "(Detalles completos en el reporte adjunto)."
    )

    # 6. Folium Map
    add_slide(
        "Resultados - Mapa Interactivo",
        "El mapa generado con Folium muestra:\n"
        "- Proximidad de sitios de lanzamiento a costas.\n"
        "- Cercanía a infraestructuras logísticas (trenes, carreteras).\n"
        "- Agrupamiento de éxitos y fallos por sitio."
    )

    # 7. Dash App
    add_slide(
        "Resultados - Dashboard Plotly",
        "La aplicación interactiva permite:\n"
        "- Filtrar por sitio de lanzamiento.\n"
        "- Seleccionar rangos de carga útil.\n"
        "- Visualizar la correlación entre éxito y carga de manera dinámica."
    )

    # 8. Machine Learning
    add_slide(
        "Resultados - Machine Learning",
        "Se evaluaron 4 modelos: Regresión Logística, SVM, Árbol de Decisión, KNN.\n"
        "Todos mostraron un rendimiento similar y alto (~83%).\n"
        "La matriz de confusión muestra una buena capacidad de predicción.",
        "assets/ml/model_comparison.png"
    )

    # 9. Conclusion
    add_slide(
        "Conclusiones",
        "- Es posible predecir el éxito con buena precisión.\n"
        "- Factores clave: Órbita, Masa de carga, Sitio de lanzamiento.\n"
        "- La reutilización es viable y predecible, apoyando el modelo de negocio de SpaceX."
    )

    # Save
    output_path = "Final_Presentation_Spacex.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
